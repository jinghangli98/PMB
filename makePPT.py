import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
import glob
import datetime as dt
from natsort import natsorted, ns

import pdb
prs = Presentation()

lyt=prs.slide_layouts[0] # choosing a slide layout
slide=prs.slides.add_slide(lyt) # adding a slide
title=slide.shapes.title # assigning a title

subtitle=slide.placeholders[1] # placeholder for subtitle
ID='ADRC_57'
date='2022.12.19-19.45.32'
CWID='CW22-44'
title.text=f'{CWID}' # title
subtitle.text=f'{date}' # subtitle

slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(100)

prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

base_path = '/Users/jinghangli/Library/CloudStorage/OneDrive-UniversityofPittsburgh/03-PMB/PMB_ADRC/'
T1_path = glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/match_T1/*.png')
T1_path = natsorted(T1_path)

T2_path = glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/match_T2/*.png')
T2_path = natsorted(T2_path)

cam_path = glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/resizedCam_adjusted/*.png')
cam_path = natsorted(cam_path)

ap_layout = glob.glob(f'{base_path}/{date}/{ID}/rembg_cam/ap/*.png')

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.add_picture(ap_layout[0], Inches(2.5), Inches(0.5))

for ind in range(len(T1_path)):
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(T1_path[ind], Inches(1), Inches(2))
    slide.shapes.add_picture(T2_path[ind], Inches(5), Inches(2))
    slide.shapes.add_picture(cam_path[ind], Inches(9), Inches(2))
# pdb.set_trace()
prs.save(f"{base_path}{date}/{CWID}.pptx") # saving file
